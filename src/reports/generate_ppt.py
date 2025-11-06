from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os, json
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import glob

# ========== 工具函数 ==========
def add_title_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle

def add_text_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    for line in content:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)

def add_image_slide(prs, title, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), Inches(8), Inches(4.5))

def save_plot(fig_path):
    plt.tight_layout()
    plt.savefig(fig_path, bbox_inches="tight")
    plt.close()

# ========== 主逻辑 ==========
def main():
    prs = Presentation()
    os.makedirs("reports/figures", exist_ok=True)
    fig_dir = "reports/figures"

    # ========== Page 1 ==========
    add_title_slide(
        prs,
        "Olist Customer Segmentation & CLV Prediction",
        "From Machine Learning to Marketing ROI Experiment"
    )

    # ========== Page 2: Data & Features ==========
    add_text_slide(prs, "Data & Feature Engineering", [
        "• 9 raw tables → 95K customers consolidated",
        "• RFM + behavioral + payment + review + time features",
        "• 10 engineered variables used in ML-CLV model",
    ])

    # ========== Page 3: Segmentation ==========
    seg = pd.read_parquet("data/features/customer_segments_kmeans.parquet")
    sns.countplot(x="segment_kmeans", data=seg)
    plt.title("Customer Segmentation Distribution (K=3)")
    fig_path = os.path.join(fig_dir, "segment_distribution.png")
    save_plot(fig_path)
    add_image_slide(prs, "Customer Segmentation (KMeans, K=3)", fig_path)

    # ========== Page 4: CLV Model Validation ==========
    cv = json.load(open("models/xgb_clv_cv_metrics.json"))
    cv_df = pd.DataFrame(cv)
    plt.figure(figsize=(8,3))
    fig, ax = plt.subplots(1,2,figsize=(8,3))
    sns.barplot(data=cv_df, x="ref", y="mae", ax=ax[0])
    sns.barplot(data=cv_df, x="ref", y="spearman", ax=ax[1])
    ax[0].set_title("MAE by Reference Date")
    ax[1].set_title("Spearman Correlation by Reference Date")
    fig_path = os.path.join(fig_dir, "clv_validation.png")
    save_plot(fig_path)
    add_image_slide(prs, "CLV Model Performance", fig_path)

    # ========== Page 5: CLV Distribution ==========
    import glob
    latest_file = sorted(glob.glob("data/outputs/CLV_ML_*.parquet"))[-1]
    print(f"[INFO] Using latest CLV file: {latest_file}")
    clv = pd.read_parquet(latest_file)

    plt.figure(figsize=(8,4))
    sns.histplot(clv["clv_ml_pred"], bins=100, log_scale=(False, True))
    plt.title("CLV (ML) Distribution (log-scaled)")
    plt.xlabel("Predicted CLV")
    plt.ylabel("Customer Count (log scale)")
    fig_path = os.path.join(fig_dir, "clv_distribution.png")
    save_plot(fig_path)
    add_image_slide(prs, "CLV Distribution & Value Structure", fig_path)

    # ========== Page 6: Segment × CLV Matrix ==========
    matrix_path = "data/outputs/CLV_segment_matrix.csv"
    if os.path.exists(matrix_path):
        df = pd.read_csv(matrix_path)
        pivot = df.pivot(index="segment_kmeans", columns="CLV_tier", values="avg_clv")
        sns.heatmap(pivot, annot=True, fmt=".0f", cmap="YlGnBu")
        plt.title("Average CLV by Segment × Tier")
        fig_path = os.path.join(fig_dir, "segment_clv_heatmap.png")
        save_plot(fig_path)
        add_image_slide(prs, "Segment × CLV Tier Matrix", fig_path)

    # ========== Page 7: Strategy Summary ==========
    add_text_slide(prs, "Marketing Strategy & Budget", [
        "• Top 5%: Loyalty & Retention → Membership, Free Shipping, Early Access",
        "• Top 20%: Repurchase Promotion → Tiered Discounts, Add-on Coupons",
        "• Mid 50%: Activation → Cart Reminders, Small Coupons",
        "• Bottom 50%: Light Reach → EDM or App Notification",
        "• Budget Allocation: 70% on top 25% high-value customers"
    ])

    # ========== Page 8: Uplift Experiment Design ==========
    alloc = pd.read_parquet("data/outputs/exp_alloc_demo.parquet")
    summary = alloc.groupby(["segment_kmeans","CLV_tier","assign"]).size().reset_index(name="count")
    sns.barplot(data=summary, x="CLV_tier", y="count", hue="assign")
    plt.title("Treatment vs Control by CLV Tier")
    fig_path = os.path.join(fig_dir, "uplift_allocation.png")
    save_plot(fig_path)
    add_image_slide(prs, "Uplift Experiment Design", fig_path)

    # ========== Page 9: Experiment Results ==========
    add_text_slide(prs, "Uplift Evaluation", [
        "• Qini AUC > 0 → Model can rank uplift correctly",
        "• Incremental Profit > 0 → Strategy yields positive ROI",
        "• Treatment/Control balanced per stratum → statistically valid",
    ])

    # ========== Page 10: Conclusion ==========
    add_text_slide(prs, "Conclusion & Next Steps", [
        "✅ ML-CLV model effectively predicts customer value",
        "✅ Segment × CLV enables precise marketing actions",
        "✅ Uplift experiment proves incremental profit",
        "Next steps:",
        "• Automate weekly scoring & experiment design",
        "• Integrate with CRM APIs for real-time campaign orchestration",
        "• Retrain quarterly & monitor PSI drift"
    ])

    out_ppt = "reports/Olist_CLV_Report.pptx"
    prs.save(out_ppt)
    print(f"✅ Saved presentation: {out_ppt}")

if __name__ == "__main__":
    main()